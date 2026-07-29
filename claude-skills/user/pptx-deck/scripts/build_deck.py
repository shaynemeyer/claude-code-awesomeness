#!/usr/bin/env python3
"""Assemble a PowerPoint deck from a JSON slide spec.

This script does NOT parse markdown itself — the calling agent reads the
source markdown, decides how to split it into slides, and writes a JSON
spec matching the schema below. This keeps slide-splitting judgment (what
belongs on one slide, where to cut a long section) with the agent instead
of a brittle markdown parser.

JSON spec schema:
{
  "title": "Deck title",
  "subtitle": "Optional subtitle",           // for the title slide only
  "slides": [
    {
      "title": "Slide title",
      "subtitle": "Optional slide subtitle",
      "bullets": [                            // optional
        {"level": 0, "text": "Top-level point"},
        {"level": 1, "text": "Sub-point"}
      ],
      "table": {                              // optional, mutually exclusive-ish with bullets
        "headers": ["Col A", "Col B"],
        "rows": [["a1", "b1"], ["a2", "b2"]]
      },
      "image": "path/to/rendered_diagram.png", // optional, e.g. a rendered mermaid diagram
      "notes": "Speaker notes / talk track for this slide"
    }
  ]
}

Usage:
  python3 build_deck.py --spec spec.json --out deck.pptx
"""
import argparse
import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROW_ALT = RGBColor(0xF2, 0xF6, 0xFA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def add_title(slide, text, subtitle=None, size=36):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.3 if subtitle else 1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.italic = True
        p2.font.color.rgb = GRAY
    line_top = Inches(1.15 if not subtitle else 1.5)
    line = slide.shapes.add_shape(1, Inches(0.6), line_top, Inches(2.0), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_bullets(slide, bullets, top=1.75, width=11.9, left=0.7, height=None):
    if height is None:
        height = 7.5 - top - 0.3
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        level = b.get("level", 0)
        text = b.get("text", "")
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = min(level, 4)
        p.font.size = Pt(22 if level == 0 else 18)
        p.font.color.rgb = DARK if level == 0 else GRAY
        p.space_after = Pt(10)


def add_table(slide, table_spec, top=1.9):
    headers = table_spec.get("headers", [])
    rows = table_spec.get("rows", [])
    n_rows = len(rows) + 1
    n_cols = len(headers)
    if n_cols == 0:
        return
    left = Inches(0.7)
    width = Inches(11.9)
    height = Inches(min(0.5 * n_rows, 5.2))
    gtable = slide.shapes.add_table(n_rows, n_cols, left, Inches(top), width, height).table

    col_widths = table_spec.get("col_widths")
    if col_widths and len(col_widths) == n_cols:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            gtable.columns[i].width = Emu(int(Inches(11.9) * w / total))

    for i, h in enumerate(headers):
        cell = gtable.cell(0, i)
        cell.text = str(h)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(15)
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gtable.cell(r, c)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(13)
            cell.text_frame.paragraphs[0].font.color.rgb = DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_ALT if r % 2 == 0 else WHITE


def add_image(slide, image_path, top=1.8, max_width=11.5, max_height=5.3):
    if not image_path or not os.path.exists(image_path):
        return
    from PIL import Image
    with Image.open(image_path) as im:
        w_px, h_px = im.size
    aspect = w_px / h_px
    width_in = max_width
    height_in = width_in / aspect
    if height_in > max_height:
        height_in = max_height
        width_in = height_in * aspect
    left_in = (13.333 - width_in) / 2
    slide.shapes.add_picture(image_path, Inches(left_in), Inches(top), width=Inches(width_in), height=Inches(height_in))


def build_title_slide(prs, title, subtitle=None, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.LEFT
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.italic = True
        p2.font.color.rgb = ACCENT
    line = slide.shapes.add_shape(1, Inches(0.8), Inches(4.6), Inches(3.0), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    set_notes(slide, notes)
    return slide


def build_content_slide(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, spec.get("title", ""), spec.get("subtitle"))
    top = 1.55 if spec.get("subtitle") else 1.75

    has_table = "table" in spec and spec["table"]
    has_image = "image" in spec and spec["image"]
    has_bullets = "bullets" in spec and spec["bullets"]

    if has_table:
        add_table(slide, spec["table"], top=top + 0.15)
    elif has_image:
        add_image(slide, spec["image"], top=top + 0.1)
    elif has_bullets:
        add_bullets(slide, spec["bullets"], top=top)

    set_notes(slide, spec.get("notes"))
    return slide


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--spec", required=True, help="path to JSON slide spec")
    parser.add_argument("--out", required=True, help="output .pptx path")
    args = parser.parse_args()

    with open(args.spec, "r", encoding="utf-8") as f:
        spec = json.load(f)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(
        prs,
        spec.get("title", "Untitled Deck"),
        spec.get("subtitle"),
        spec.get("title_notes"),
    )

    for slide_spec in spec.get("slides", []):
        build_content_slide(prs, slide_spec)

    prs.save(args.out)
    print("wrote %s (%d slides)" % (args.out, len(prs.slides)))


if __name__ == "__main__":
    main()
